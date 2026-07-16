"""Export project tutorial to PowerPoint (no charts)."""
from pptx import Presentation
from pptx.util import Inches, Pt
from schema import ProjectTutorial


def _title_slide(prs, title, subtitle=""):
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:100]
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle[:300]


def _bullets(prs, title, lines):
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title[:100]
    body = slide.placeholders[1].text_frame
    lines = [str(x) for x in (lines or [])][:40]
    if not lines:
        body.text = ""
        return
    body.text = lines[0][:900]
    for line in lines[1:]:
        p = body.add_paragraph()
        p.text = line[:900]


def export_tutorial_to_pptx(tutorial: ProjectTutorial, output_path: str):
    prs = Presentation()
    _title_slide(
        prs,
        tutorial.title,
        f"{tutorial.tagline}\n{tutorial.difficulty} · {tutorial.size} · ~{tutorial.estimated_hours}h",
    )
    _bullets(prs, "What we're building", [tutorial.problem_statement, tutorial.end_state])
    _bullets(prs, "Stack", [f"{s.name} {s.version} — {s.purpose}" for s in tutorial.stack])
    _bullets(prs, "Repo layout", tutorial.repo_layout)
    _bullets(prs, "Prerequisites", tutorial.prerequisites or ["None"])

    if tutorial.intro_transcript:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
        box.text_frame.text = "Intro"
        s.notes_slide.notes_text_frame.text = tutorial.intro_transcript[:20000]

    for step in tutorial.steps:
        blank = prs.slide_layouts[6]
        s = prs.slides.add_slide(blank)
        tb = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        tf = tb.text_frame
        tf.text = f"Step {step.step_number}: {step.title}"
        tf.paragraphs[0].font.size = Pt(24)
        tf.paragraphs[0].font.bold = True

        bb = s.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(9), Inches(2.2))
        btf = bb.text_frame
        btf.word_wrap = True
        bullets = step.bullets or [step.goal]
        btf.text = "• " + bullets[0]
        for b in bullets[1:]:
            p = btf.add_paragraph()
            p.text = "• " + b
            p.font.size = Pt(14)
        btf.paragraphs[0].font.size = Pt(14)

        note = step.speaker_notes or ""
        if step.code:
            cb = s.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(3.2))
            ctf = cb.text_frame
            ctf.word_wrap = True
            ctf.text = f"# {step.code.file_path}\n{step.code.code}"[:3500]
            for para in ctf.paragraphs:
                para.font.name = "Consolas"
                para.font.size = Pt(10)
            note += f"\n\n[CODE {step.code.file_path}]\n{step.code.code}"
        s.notes_slide.notes_text_frame.text = note[:20000]

    if tutorial.outro_transcript:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.5))
        box.text_frame.text = "Outro"
        s.notes_slide.notes_text_frame.text = tutorial.outro_transcript[:20000]

    if tutorial.final_project_code:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        box.text_frame.text = "Final project code"
        cb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
        cb.text_frame.word_wrap = True
        cb.text_frame.text = tutorial.final_project_code[:5000]

    _bullets(prs, "How to run", [tutorial.how_to_run])

    for ch in tutorial.challenges or []:
        _bullets(prs, f"Challenge: {ch.title}", [ch.description, ch.build_process])
        s = prs.slides.add_slide(prs.slide_layouts[6])
        box = s.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.5))
        box.text_frame.text = f"Solution: {ch.title}"
        cb = s.shapes.add_textbox(Inches(0.5), Inches(1.0), Inches(9), Inches(5.5))
        cb.text_frame.word_wrap = True
        cb.text_frame.text = ch.solution_code[:5000]
        s.notes_slide.notes_text_frame.text = ch.solution_walkthrough[:20000]

    prs.save(output_path)
    return output_path


# legacy name
def export_course_to_pptx(course, output_path: str = "course.pptx"):
    return export_tutorial_to_pptx(course, output_path)